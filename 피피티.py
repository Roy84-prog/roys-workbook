import os
import json
import glob
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE


# ---------------------------------------------------------
# 1. 텍스트 파싱 및 서식 적용 함수
# ---------------------------------------------------------
def add_formatted_text_to_paragraph(paragraph, text):
    pattern = re.compile(r'\[([A-Z])\](.*?)\[/\1\]|《(.*?)》(.*?)《/\3》')
    last_idx = 0
    for match in pattern.finditer(text):
        normal_text = text[last_idx:match.start()]
        if normal_text:
            run = paragraph.add_run()
            run.text = normal_text
            run.font.size = Pt(24)
            run.font.name = '맑은 고딕'

        bracket_tag = match.group(1)
        bracket_text = match.group(2)
        angle_tag = match.group(3)
        angle_text = match.group(4)

        run = paragraph.add_run()
        run.font.size = Pt(24)
        run.font.name = '맑은 고딕'
        run.font.bold = True

        if bracket_tag:
            run.text = bracket_text
            if bracket_tag == 'S':
                run.font.color.rgb = RGBColor(0, 112, 192)
            elif bracket_tag == 'V':
                run.font.color.rgb = RGBColor(255, 0, 0)
            elif bracket_tag == 'O':
                run.font.color.rgb = RGBColor(0, 176, 80)
            elif bracket_tag == 'C':
                run.font.color.rgb = RGBColor(112, 48, 160)
        elif angle_tag:
            run.text = angle_text
            if angle_tag in ['clue', 'feature']:
                run.font.color.rgb = RGBColor(255, 102, 0)
                if angle_tag == 'clue': run.font.underline = True
            elif angle_tag == 'subject':
                run.font.color.rgb = RGBColor(0, 112, 192)
                run.font.italic = True
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

        last_idx = match.end()

    remaining_text = text[last_idx:]
    if remaining_text:
        run = paragraph.add_run()
        run.text = remaining_text
        run.font.size = Pt(24)
        run.font.name = '맑은 고딕'


# ---------------------------------------------------------
# 2. 단일 PPT 생성 엔진
# ---------------------------------------------------------
def create_single_ppt(data, output_filepath):
    prs = Presentation()

    # --- 표지 슬라이드 ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])

    # "울산고1 내신 기말고사" 타이틀
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = "울산고1 내신 기말고사"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.name = '맑은 고딕'

    # "Roy (심성훈) 선생님" - 지시하신 대로 더 낮고, 작게, 같은 폰트로 설정
    sub_box = title_slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    sub_p = sub_box.text_frame.add_paragraph()
    sub_p.text = "Roy (심성훈) 선생님"
    sub_p.font.size = Pt(28)
    sub_p.font.name = '맑은 고딕'

    # 지문 정보 출력 부분
    info_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    info_p = info_box.text_frame.add_paragraph()
    source = data.get('meta_info', {}).get('source_origin', '출처 미상')
    topic = data.get('topic_data', {})
    info_p.text = f"출처: {source}\n{topic.get('keywords', '')}\n\n{topic.get('summary', '')}"
    info_p.font.size = Pt(18)

    # --- 문장 분석 슬라이드 ---
    for sentence in data.get("sentence_analysis", []):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Sentence {sentence.get('num', '')} : {sentence.get('summary_mark', '')}"

        # [영어 텍스트 박스]
        eng_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(2.2))
        eng_tf = eng_box.text_frame
        eng_tf.word_wrap = True
        eng_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        eng_paragraph = eng_tf.add_paragraph()
        eng_paragraph.line_spacing = 1.3
        add_formatted_text_to_paragraph(eng_paragraph, sentence.get("eng_analyzed", ""))

        # [한글 텍스트 박스]
        kor_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9.0), Inches(1.8))
        kor_tf = kor_box.text_frame
        kor_tf.word_wrap = True
        kor_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        kor_paragraph = kor_tf.add_paragraph()
        kor_paragraph.line_spacing = 1.3
        kor_paragraph.text = sentence.get("kor_translation", "")
        kor_paragraph.font.size = Pt(22)
        kor_paragraph.font.color.rgb = RGBColor(100, 100, 100)
        kor_paragraph.font.name = '맑은 고딕'

        # [발표자 노트]
        notes_frame = slide.notes_slide.notes_text_frame
        notes = f"★ 쉬운 풀이: {sentence.get('easy_exp', '')}\n\n★ 문맥 어휘: {sentence.get('context_meaning', '')}\n\n★ 구문 팁:\n"
        for tip in sentence.get("syntax_tip", []):
            notes += f"{tip.get('tag', '')} {tip.get('target', '')}: {tip.get('explanation', '')}\n"
        notes_frame.text = notes

    prs.save(output_filepath)


# ---------------------------------------------------------
# 3. 바탕화면 '3' 폴더 일괄 처리 컨트롤러
# ---------------------------------------------------------
def process_all_files_in_folder():
    # 맥북 바탕화면의 '3' 폴더 경로로 수정 완료
    folder_path = os.path.expanduser("~/Desktop/3")

    if not os.path.exists(folder_path):
        print(f"오류: '{folder_path}' 경로를 찾을 수 없습니다. 바탕화면에 '3' 폴더가 있는지 확인해 주세요.")
        return

    txt_files = glob.glob(os.path.join(folder_path, "*.txt"))

    if not txt_files:
        print(f"'{folder_path}' 폴더 안에 .txt 파일이 없습니다.")
        return

    print(f"총 {len(txt_files)}개의 txt 파일을 발견했습니다. 일괄 변환을 시작합니다...\n")

    success_count = 0
    for txt_file in txt_files:
        filename = os.path.basename(txt_file)
        base_name = os.path.splitext(filename)[0]
        output_filepath = os.path.join(folder_path, f"{base_name}.pptx")

        try:
            with open(txt_file, 'r', encoding='utf-8') as f:
                content = f.read()
                data = json.loads(content)

            create_single_ppt(data, output_filepath)
            print(f"✅ 성공: {filename} ➔ {base_name}.pptx")
            success_count += 1

        except json.JSONDecodeError:
            print(f"❌ 실패: {filename} (JSON 형식이 올바르지 않습니다. AI가 생성한 텍스트 구조를 확인하세요.)")
        except Exception as e:
            print(f"❌ 에러 발생 ({filename}): {str(e)}")

    print(f"\n작업 완료! 총 {success_count}/{len(txt_files)}개의 PPT가 성공적으로 생성되었습니다.")


if __name__ == "__main__":
    process_all_files_in_folder()